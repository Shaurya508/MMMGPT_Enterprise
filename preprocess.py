import io
import os
import re
import time
import sys
from pptx import Presentation
from PIL import Image
import numpy as np
import easyocr
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_ollama import OllamaEmbeddings
import json 
from typing import Iterator
from langchain_core.document_loaders import BaseLoader
from langchain_core.documents import Document as LCDocument
from docling.document_converter import DocumentConverter
# import google.generativeai as genai
from paddleocr import PaddleOCR

with open('config.json', 'r') as f:
    config = json.load(f)

# Set environment variables from config file
for key, value in config["ENV_VARS"].items():
    os.environ[key] = value

# Use config for embedding model and chat model
EMBEDDING_MODEL = config["EMBEDDING_MODEL"]
CHAT_MODEL = config["CHAT_MODEL"]

CHUNK_SIZE = config["CHUNK_SIZE"]
CHUNK_OVERLAP = config["CHUNK_OVERLAP"]
DB_NAME = config["DB_NAME"]

# Original code

# # PaddleOCR-based function
# from paddleocr import PaddleOCR
# from PIL import Image

# from pptx import Presentation
# from paddleocr import PaddleOCR
# from PIL import Image
# import io
# import os

# # Directory to save extracted images
# EXTRACTED_IMG_DIR = "extracted_images"
# os.makedirs(EXTRACTED_IMG_DIR, exist_ok=True)

# def ocr_with_paddle(image_path):
#     """
#     Extract text from an image using PaddleOCR.

#     Args:
#         image_path (str): Path to the image file.

#     Returns:
#         str: Extracted text from the image.
#     """
#     # Initialize PaddleOCR
#     ocr = PaddleOCR(lang='en', use_angle_cls=True)

#     # Perform OCR on the image
#     result = ocr.ocr(image_path)

#     # Extract the text from the OCR result
#     extracted_text = ''
#     for line in result[0]:
#         extracted_text += ' ' + line[1][0]

#     return extracted_text.strip()


# def extract_text_from_image(shape, slide_number, image_number):
#     """
#     Extract text from images in PPT slides using PaddleOCR.

#     Args:
#         shape: The shape object containing the image in the PPT.
#         slide_number: The slide number of the image.
#         image_number: The image number on the slide.

#     Returns:
#         A formatted string with extracted text wrapped in flags.
#     """
#     try:
#         # Extract image bytes and convert to PIL Image
#         image_bytes = shape.image.blob
#         image_stream = io.BytesIO(image_bytes)
#         img = Image.open(image_stream)

#         # Save the image locally
#         image_path = os.path.join(EXTRACTED_IMG_DIR, f"slide_{slide_number}_image_{image_number}.png")
#         img.save(image_path)

#         # Use PaddleOCR to extract text from the image
#         extracted_text = ocr_with_paddle(image_path)

#         # Add flags around the extracted text
#         image_opening_flag = f"<|S{slide_number}I{image_number}|>"
#         image_closing_flag = f"</|S{slide_number}I{image_number}|>"

#         return f"{image_opening_flag}\n{extracted_text}\n{image_closing_flag}"
#     except Exception as e:
#         print(f"Error processing image on slide {slide_number}, image {image_number}: {e}")
#         return f"<|S{slide_number}I{image_number}|>\nError extracting text\n</|S{slide_number}I{image_number}|>"


# Modified code

import os
import io
import json
from PIL import Image
from paddleocr import PaddleOCR

# Directory to save extracted images
EXTRACTED_IMG_DIR = "extracted_images"
os.makedirs(EXTRACTED_IMG_DIR, exist_ok=True)

def ocr_with_paddle(image_path):
    """
    Extract text and bounding box data from an image using PaddleOCR.

    Args:
        image_path (str): Path to the image file.

    Returns:
        list: Extracted text lines with bounding box coordinates.
    """
    # Initialize PaddleOCR
    ocr = PaddleOCR(lang='en', use_angle_cls=True)

    # Perform OCR on the image
    result = ocr.ocr(image_path)

    # Return structured data (text + bounding boxes)
    return result[0]

def parse_table_data(ocr_results):
    """
    Parse OCR results to organize text into a table-like structure.

    Args:
        ocr_results (list): List of OCR results with text and bounding boxes.

    Returns:
        list: Extracted table rows as a list of dictionaries.
    """
    table_data = []
    for line in ocr_results:
        # Extract the bounding box and text
        bbox = line[0]
        text = line[1][0]
        
        # Append to table data
        table_data.append({"text": text, "bbox": bbox})

    return table_data

def format_table_data(parsed_data):
    """
    Convert parsed data into a JSON-like table format.

    Args:
        parsed_data (list): Parsed OCR data.

    Returns:
        str: JSON representation of the table.
    """
    # Convert to structured JSON format
    table = []
    for item in parsed_data:
        table.append({"text": item["text"], "bbox": item["bbox"]})
    
    # Return as JSON string
    return json.dumps(table, indent=4)

def extract_text_from_image(shape, slide_number, image_number):
    """
    Extract text from images in PPT slides using PaddleOCR.

    Args:
        shape: The shape object containing the image in the PPT.
        slide_number: The slide number of the image.
        image_number: The image number on the slide.

    Returns:
        A formatted string with extracted text wrapped in flags.
    """
    try:
        # Extract image bytes and convert to PIL Image
        image_bytes = shape.image.blob
        image_stream = io.BytesIO(image_bytes)
        img = Image.open(image_stream)

        # Save the image locally
        image_path = os.path.join(EXTRACTED_IMG_DIR, f"slide_{slide_number}_image_{image_number}.png")
        img.save(image_path)

        # Use PaddleOCR to extract text and bounding box data
        ocr_results = ocr_with_paddle(image_path)

        # Parse and format the extracted data
        parsed_data = parse_table_data(ocr_results)
        formatted_text = format_table_data(parsed_data)

        # Add flags around the extracted text
        image_opening_flag = f"<|S{slide_number}I{image_number}|>"
        image_closing_flag = f"</|S{slide_number}I{image_number}|>"

        return f"{image_opening_flag}\n{formatted_text}\n{image_closing_flag}"
    except Exception as e:
        print(f"Error processing image on slide {slide_number}, image {image_number}: {e}")
        return f"<|S{slide_number}I{image_number}|>\nError extracting text\n</|S{slide_number}I{image_number}|>"

# # Example: Process the uploaded image
# image_path = "/mnt/data/image.png"  # Replace with actual path
# ocr_results = ocr_with_paddle(image_path)
# parsed_data = parse_table_data(ocr_results)
# formatted_data = format_table_data(parsed_data)

# print(formatted_data)  # Output as JSON


def extract_title(shape, slide_number, chart_count):
    """Extract chart title and save chart image from the slide"""
    chart_title = ""
    if shape.chart.has_title:
        chart_title = shape.chart.chart_title.text

    title_opening_flag = f"<|S{slide_number}C{chart_count}|>"
    title_closing_flag = f"</|S{slide_number}C{chart_count}|>"
    return f"{title_opening_flag}\n{chart_title}\n{title_closing_flag}"


def get_ppt_data(file_name):
    """Extract text, tables, and images from PPT"""
    prs = Presentation(file_name)
    extracted_text = ""
    total_slides = len(prs.slides)

    if total_slides == 0:
        print("No slides found in the presentation.")
        return extracted_text

    processed_slides = 0

    for slide_number, slide in enumerate(prs.slides, start=1):
        image_count = 0
        chart_count = 0

        for shape in slide.shapes:
            try:
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        extracted_text += str(row_data) + '\n'

                if hasattr(shape, "text"):
                    text = shape.text.replace('\v', '\n').replace('\x0b', '\n')
                    extracted_text += text + '\n'

                if hasattr(shape, "image"):
                    image_count += 1
                    extracted_text += extract_text_from_image(shape, slide_number, image_count) + '\n'

                if shape.has_chart:
                    chart_count += 1
                    extracted_text += extract_title(shape, slide_number, chart_count) + '\n'
            except Exception as e:
                print(f"Error processing shape on slide {slide_number}: {e}")

        # Update progress after processing each slide
        processed_slides += 1
        progress = 10 + int((processed_slides / total_slides) * 50)  # First 50% of progress
        print(f"PROGRESS:{progress}", flush=True)

    return extracted_text

# class DoclingPDFLoader(BaseLoader):

#     def __init__(self, file_path: str | list[str]) -> None:
#         self._file_paths = file_path if isinstance(file_path, list) else [file_path]
#         self._converter = DocumentConverter()

#     def lazy_load(self) -> Iterator[LCDocument]:
#         for source in self._file_paths:
#             dl_doc = self._converter.convert(source).document
#             text = dl_doc.export_to_markdown()
#             yield LCDocument(page_content=text)

# def extract_text_from_image(shape, reader, slide_number, image_number):
#     """Extract text from images in PPT slides"""
#     image_bytes = shape.image.blob
#     image_stream = io.BytesIO(image_bytes)
#     img = Image.open(image_stream)

#     image_path = os.path.join(EXTRACTED_IMG_DIR, f"slide_{slide_number}_image_{image_number}.png")
#     img.save(image_path)
#     img_np = np.array(img.convert('L'))
#     result = reader.readtext(img_np, detail=1)

#     rows = []
#     current_row = []
#     previous_y = None
#     tolerance = 15

#     for entry in result:
#         bbox, text, confidence = entry
#         x_min, y_min = bbox[0]

#         if previous_y is None or abs(y_min - previous_y) <= tolerance:
#             current_row.append(text)
#         else:
#             rows.append(current_row)
#             current_row = [text]
#         previous_y = y_min

#     if current_row:
#         rows.append(current_row)

#     table_str = ""
#     for row in rows:
#         table_str += '\t'.join(row) + '\n'

#     image_opening_flag = f"<|S{slide_number}I{image_number}|>"
#     image_closing_flag = f"</|S{slide_number}I{image_number}|>"
#     return f"{image_opening_flag}\n{table_str}\n{image_closing_flag}"

# def extract_title(shape, slide_number, chart_count):
#     """Extract chart title and save chart image from the slide"""
#     chart_title = ""
#     if shape.chart.has_title:
#         chart_title = shape.chart.chart_title.text_frame.text

#     title_opening_flag = f"<|S{slide_number}C{chart_count}|>"
#     title_closing_flag = f"</|S{slide_number}C{chart_count}|>"
#     return f"{title_opening_flag}\n{chart_title}\n{title_closing_flag}"

# def get_ppt_data(file_name):
#     """Extract text, tables, and images from PPT"""
#     prs = Presentation(file_name)
#     extracted_text = ""
#     reader = easyocr.Reader(['en'])
#     total_slides = len(prs.slides)
#     processed_slides = 0

#     # Use DoclingPDFLoader to extract the main textual content
#     loader = DoclingPDFLoader(file_path=file_name)
#     docs = loader.load()
#     extracted_text += docs[0].page_content  # Assuming single document content

#     for slide_number, slide in enumerate(prs.slides, start=1):
#         image_count = 0
#         chart_count = 0

#         for shape in slide.shapes:
#             if hasattr(shape, "image"):
#                 image_count += 1
#                 extracted_text += extract_text_from_image(shape, reader, slide_number, image_count) + '\n'

#             if shape.has_chart:
#                 chart_count += 1
#                 extracted_text += extract_title(shape, slide_number, chart_count) + '\n'

#         # Update progress after processing each slide
#         processed_slides += 1
#         progress = 10 + int((processed_slides / total_slides) * 50)  # First 50% of progress
#         print(f"PROGRESS:{progress}", flush=True)

#     return extracted_text



def get_text_chunks_with_flags(content):
    """Chunk text with metadata flags"""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE, 
        chunk_overlap=CHUNK_OVERLAP
    )
    chunks_with_metadata = []
    current_metadata = set()

    image_flag_pattern = r"<\|S(\d+)I(\d+)\|>"
    chart_flag_pattern = r"<\|S(\d+)C(\d+)\|>"
    image_closing_pattern = r"</\|S(\d+)I(\d+)\|>"
    chart_closing_pattern = r"</\|S(\d+)C(\d+)\|>"

    text_chunks = text_splitter.split_text(content)
    
    for chunk in text_chunks:
        chunk_metadata = set()

        # Process flags and maintain metadata
        for pattern, closing_pattern, prefix in [
            (image_flag_pattern, image_closing_pattern, "slide_%s_image_%s.png"),
            (chart_flag_pattern, chart_closing_pattern, "slide_%s_chart_%s.png")
        ]:
            opening_flags = re.findall(pattern, chunk)
            for match in opening_flags:
                file_path = prefix % match
                chunk_metadata.add(file_path)
                current_metadata.add(file_path)

            closing_flags = re.findall(closing_pattern, chunk)
            for match in closing_flags:
                file_path = prefix % match
                chunk_metadata.add(file_path)
                current_metadata.discard(file_path)

        chunk_metadata.update(current_metadata)

        # Clean chunk text
        chunk_clean = re.sub(image_flag_pattern, '', chunk)
        chunk_clean = re.sub(image_closing_pattern, '', chunk_clean)
        chunk_clean = re.sub(chart_flag_pattern, '', chunk_clean)
        chunk_clean = re.sub(chart_closing_pattern, '', chunk_clean)

        chunks_with_metadata.append({
            "text": chunk_clean.strip(),
            "metadata": {"image_paths": list(chunk_metadata)}
        })

    return chunks_with_metadata

if __name__ == "__main__":
    start_time = time.time()
    
    if len(sys.argv) < 4:
        print("Error: Too few arguments. Usage: python preprocess.py <ppt_file_path> <ppt_dir> <main_img_dir>")
        sys.exit(1)
    elif len(sys.argv) > 4:
        print("Error: Too many arguments. Usage: python preprocess.py <ppt_file_path> <ppt_dir> <main_img_dir>")
        sys.exit(1)

    ppt_file_name = sys.argv[1]
    PPT_DIR = sys.argv[2]
    MAIN_IMG_DIR = sys.argv[3]

    FULL_SLIDE_IMG_DIR = os.path.join(MAIN_IMG_DIR, "Full Slide Images")
    EXTRACTED_IMG_DIR = os.path.join(MAIN_IMG_DIR, "Extracted Images")
    EXTRACTED_CHART_DIR = os.path.join(MAIN_IMG_DIR, "Extracted Charts")
    
    # Initialize embedding model
    embeddings = OllamaEmbeddings(model=EMBEDDING_MODEL)

    print("PROGRESS:10", flush=True)
    
    # Extract and process PPT data
    extracted_text = get_ppt_data(ppt_file_name)

    # Update progress to 60% after extracting PPT data
    print("PROGRESS:60", flush=True)

    chunks_with_metadata = get_text_chunks_with_flags(extracted_text)

    # Update progress to 70% after chunking
    print("PROGRESS:70", flush=True)
    
    # Create FAISS database
    texts = [item['text'] for item in chunks_with_metadata]
    metadatas = [item['metadata'] for item in chunks_with_metadata]
    vector_store = FAISS.from_texts(texts, embeddings, metadatas=metadatas)

    # Update progress to 90% after creating vector store
    print("PROGRESS:90", flush=True)

    db_save_path = os.path.join(PPT_DIR, DB_NAME)
    vector_store.save_local(db_save_path)

    # Update progress to 100% after saving the database
    print("PROGRESS:100", flush=True)

    #print(f"FAISS database created and saved as '{DB_NAME}' in the directory {PPT_DIR}")
    #print(f"Overall Time Taken = {time.time() - start_time:.2f} seconds")

    # Print a completion message
    print("Processing completed successfully.", flush=True)
